# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 색상 정의 (Premium Dark Mode Palette)
BG_COLOR = RGBColor(15, 23, 42)        # Slate 900 (매우 어두운 네이비 그레이)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Slate 800 (카드 배경)
TEXT_WHITE = RGBColor(255, 255, 255)   # 타이틀 흰색
TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400 (설명/본문 회색)
ACCENT_TEAL = RGBColor(20, 184, 166)   # Teal 500 (핵심 하이라이트)
ACCENT_BLUE = RGBColor(59, 130, 246)   # Blue 500 (서브 하이라이트)

FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Arial"

def create_slide_with_bg(prs):
    """어두운 배경이 적용된 슬라이드 생성"""
    blank_layout = prs.slide_layouts[6] # 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)
    
    # 배경색 지정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    return slide

def add_header(slide, title_text, category_text="PORTFOLIO PRESENTATION"):
    """상단 공통 헤더 추가"""
    # 카테고리 텍스트
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_BODY
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_TEAL
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def add_card(slide, left, top, width, height, title_text, text_lines):
    """SaaS 스타일 내용 카드 추가"""
    # 카드 모양 (둥근 직사각형)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG_COLOR
    shape.line.color.rgb = RGBColor(71, 85, 105) # Slate 600 경계선
    shape.line.width = Pt(1)
    
    # 카드 텍스트 프레임 생성
    tb = slide.shapes.add_textbox(
        Inches(left + 0.3), Inches(top + 0.3), 
        Inches(width - 0.6), Inches(height - 0.6)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # 카드 타이틀
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = ACCENT_TEAL
    p_title.space_after = Pt(14)
    
    # 카드 본문 라인 추가
    for line in text_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
        # 특정 키워드가 포함될 경우 폰트 색상 및 굵기 하이라이트
        if "★" in line or "**" in line:
            # 마커 기호 제거 후 가공
            clean_text = line.replace("★", "").replace("**", "")
            p.text = clean_text
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
        elif "핵심 역할" in line or "분석 기법" in line or "성과" in line:
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE

def set_presenter_notes(slide, notes_text):
    """슬라이드 노트(발표자 폰 메모) 설정"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

def main():
    prs = Presentation()
    # 16:9 와이드스크린 치수 설정 (13.33인치 x 7.5인치)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # =========================================================================
    # SLIDE 1: Title Slide (표지)
    # =========================================================================
    slide1 = create_slide_with_bg(prs)
    
    # 좌측 장식 바 (Teal)
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_TEAL
    accent_bar.line.fill.background()
    
    # 타이틀 텍스트 박스
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # 카테고리 레이블
    p0 = tf.paragraphs[0]
    p0.text = "DATA ANALYTICS & DECISION PORTFOLIO"
    p0.font.name = FONT_BODY
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_TEAL
    p0.space_after = Pt(20)
    
    # 메인 타이틀
    p1 = tf.add_paragraph()
    p1.text = "비정형 iGaming 로그 분석을 통한\n플레이어 성향 분류 및 의사결정 시뮬레이터 구축"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(14)
    
    # 서브타이틀
    p2 = tf.add_paragraph()
    p2.text = "비정형 데이터 파이프라인(ETL)부터 통계 분석, 머신러닝, 실시간 대시보드 웹 배포까지"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(40)
    
    # 메타정보
    p3 = tf.add_paragraph()
    p3.text = "SK Planet 데이터 직무 PT 면접 발표자료 (5~10분 분량)\n발표자: 지원자"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT_BLUE
    p3.font.bold = True
    
    # 슬라이드 1 노트 설정
    set_presenter_notes(slide1, """[슬라이드 1 발표 대본]
"안녕하십니까, SK Planet 데이터 직무 지원자 [지원자 이름]입니다.
저는 오늘 '비정형 iGaming 로그 분석을 통한 플레이어 성향 분류 및 의사결정 시뮬레이터 구축' 프로젝트를 주제로 발표를 진행하겠습니다.

본 프로젝트는 게임 시스템이 남긴 정제되지 않은 문자열 로그에서 출발하여, 핵심 지표를 계산하고, 통계 분석과 머신러닝을 적용해 유저의 의사결정을 돕는 웹 플랫폼까지 완성해 낸 End-to-End 데이터 파이프라인 프로젝트입니다. 발표를 시작하겠습니다."

[용어 리마인더]
- 비정형 데이터: 컴퓨터가 바로 계산할 수 없는 자유 줄글 형식의 로그 텍스트 파일.""")
    
    # =========================================================================
    # SLIDE 2: Background & Objectives (배경 및 목적)
    # =========================================================================
    slide2 = create_slide_with_bg(prs)
    add_header(slide2, "01. 프로젝트 배경 및 비즈니스 목적", "BACKGROUND & OBJECTIVES")
    
    add_card(
        slide2, 0.8, 1.8, 5.6, 4.8, 
        "비즈니스 배경 (Business Context)",
        [
            "■ iGaming/포커 플랫폼의 최우선 과제",
            "  - 고객 리텐션(Retention) 극대화 및 유저 파산/이탈 방어",
            "  - 올바른 플레이 가이드를 제시하지 못할 경우 유저 이탈로 직결",
            "",
            "■ 데이터 분석의 요구 사항",
            "  - 플레이어의 행동 로그에서 플레이 스타일 지표를 수학적으로 도출",
            "  - ★플레이 성향 지표와 수익성(Chips Won) 간의 통계적 인과관계 증명 필요",
            "  - 이를 통해 개인 맞춤형 피드백 및 의사결정 보조 시스템의 정당성 확보"
        ]
    )
    
    add_card(
        slide2, 6.9, 1.8, 5.6, 4.8, 
        "핵심 분석 및 해결 목적 (Core Objectives)",
        [
            "■ 1단계: 비정형 데이터의 구조화 및 파이프라인 수립",
            "  - 비형식적인 텍스트 로그 데이터를 관계형 DB 스키마로 가공 적재",
            "",
            "■ 2단계: 통계 분석 및 머신러닝 모델링",
            "  - 진입 전략에 따른 수익 차이 T-검정 및 플레이 성향별 군집 정의",
            "  - 유저 지표 기준 흑자/적자 분류 성능 80% 이상 확보",
            "",
            "■ 3단계: 분석 결과의 플랫폼 서비스화 (Streamlit)",
            "  - 분석 결과물인 시뮬레이터를 대화형 플랫폼으로 배포하여 유저가 스스로 플레이 전략 변경 시의 미래 자산 추이를 모의실험하도록 유도"
        ]
    )
    
    # 슬라이드 2 노트 설정
    set_presenter_notes(slide2, """[슬라이드 2 발표 대본]
"첫 번째로 프로젝트의 배경과 비즈니스 목적에 대해 말씀드리겠습니다.

포커나 카지노 같은 베팅 플랫폼 서비스에서 가장 해결해야 할 심각한 문제는 '유저의 이탈, 즉 Churn'입니다. 대다수의 유저들은 올바른 플레이 습관을 알지 못해 무분별하게 칩을 잃고 파산하게 되며, 이는 서비스 탈퇴로 이어집니다.

따라서 본 프로젝트의 궁극적인 비즈니스 목적은 유저의 베팅 행동 지표들을 추출하여 실제 수익성과의 통계적 인과관계를 밝히고, 올바른 베팅 가이드를 제시하여 이탈을 방지하는 것입니다.

이를 위해 저는 데이터를 단순히 들여다보고 보고서를 쓰는 데 그치지 않고, 비정형 로그 정제, 통계적 가설 검정과 머신러닝, 그리고 유저가 직접 내 미래 수익을 시뮬레이션할 수 있는 Streamlit 대시보드 웹앱 서비스까지 구축하는 해결 로드맵을 수립했습니다."

[용어 리마인더]
- Churn(이탈): 포커 게임에서는 올바른 습관을 가지지 못해 파산하고 아예 앱을 삭제하는 최악의 유저 행동.
- 리텐션(Retention): 고객이 앱을 지우지 않고 계속 이용하는 유지율.""")

    # =========================================================================
    # SLIDE 3: Data Specification & Pipeline (데이터와 역할)
    # =========================================================================
    slide3 = create_slide_with_bg(prs)
    add_header(slide3, "02. 비정형 로그 정제 및 데이터 파이프라인 구축", "DATA & ETL PIPELINE")
    
    add_card(
        slide3, 0.8, 1.8, 5.6, 4.8, 
        "활용 데이터 스펙 (Data Specification)",
        [
            "■ 원천 데이터 세트",
            "  - 140개 비정형 텍스트 로그 파일 (Hand History)",
            "  - 포커 테이블 내 모든 실시간 이벤트를 시간순으로 기록한 로우 로그",
            "",
            "■ 데이터 볼륨",
            "  - ★총 7,993개 핸드 (Games), 932명 플레이어 데이터",
            "  - ★10만 행 이상의 플레이어별 상세 액션 로그 분석",
            "    (Pre-flop, Flop, Turn, River 단계별 Bet, Call, Raise, Fold)",
            "",
            "■ 핵심 챌린지",
            "  - 비형식적인 문자열 데이터에서 데이터 일관성을 해치지 않고 플레이어 포지션, 배팅 칩스, 액션 흐름을 누락 없이 파싱하는 프로세스 설계"
        ]
    )
    
    add_card(
        slide3, 6.9, 1.8, 5.6, 4.8, 
        "본인이 수행한 핵심 역할 (Role & Contribution)",
        [
            "■ 1. Regex 기반 고속 파서 자체 설계",
            "  - Python 정규표현식(Regex)을 복합 활용하여 비정형 포커 액션 데이터를 시간 단위, 플레이어 단위로 완벽 분리 및 추출 성공",
            "",
            "■ 2. Star Schema 데이터 웨어하우스 모델링",
            "  - 분석 쿼리 성능과 일관성을 높이기 위해 다차원 스타 스키마 설계",
            "  - Fact Table: `hands`, `actions` (게임 결과 및 배팅 내역)",
            "  - Dimension Table: `players`, `tournaments` (유저 및 토너먼트 유형 정보)",
            "",
            "■ 3. DB 마이그레이션 호환성 레이어 구축",
            "  - SQLite 파일 DB와 PostgreSQL 엔터프라이즈 RDBMS의 접속을 동적 매핑하여 단일 환경 변수 설정으로 ETL 적재 스위칭 가능하도록 설계"
        ]
    )
    
    # 슬라이드 3 노트 설정
    set_presenter_notes(slide3, """[슬라이드 3 발표 대본]
"두 번째로 활용한 데이터 스펙과 파이프라인 구축 과정에서 제가 수행한 역할입니다.

원천 데이터는 140개의 비정형 포커 게임 텍스트 로그 파일입니다. 약 7,993개의 게임 판수와 932명의 플레이어가 남긴 10만 행 이상의 상세 배팅 로그로 구성되어 있었습니다.

저는 이를 구조화하기 위해 다음과 같은 세 가지 기술적 기여를 달성했습니다.

첫째, Python 정규표현식(Regex) 엔진을 자체 설계하여, 텍스트 줄글에서 배팅 단계, 플레이어명, 칩 수량, 베팅 종류 등을 정확하게 파싱하여 추출했습니다.

둘째, 데이터의 다차원 분석과 쿼리 성능 향상을 위해 스타 스키마(Star Schema) 모델링을 수행했습니다. 실제 배팅 액션을 담은 Fact 테이블과, 플레이어 및 토너먼트 속성을 담은 Dimension 테이블을 분리하여 데이터 웨어하우스의 표준 구조를 확립했습니다.

마지막으로, 로컬 분석용 가벼운 SQLite 데이터베이스와 실제 배포 서버용 PostgreSQL 데이터베이스를 단 하나의 환경 변수로 넘나들 수 있도록 이중 DB 호환성 코드를 구축하여 이식성을 극대화했습니다."

[용어 리마인더]
- 정규표현식 (Regex): 수많은 텍스트 중 '원하는 규칙의 글자'만 칼같이 도려내는 컴퓨터 필터 규칙.
- 스타 스키마 (Star Schema): 마트 진열처럼 핵심 기록(Fact)을 중심에 두고 주변에 설명 정보(Dimension)를 별 모양처럼 연결해 검색 효율을 높인 것.""")

    # =========================================================================
    # SLIDE 4: Core Statistics & ML (가설 검정 및 모델링)
    # =========================================================================
    slide4 = create_slide_with_bg(prs)
    add_header(slide4, "03. 통계적 가설 검정 및 머신러닝 분석", "STATISTICS & MACHINE LEARNING")
    
    add_card(
        slide4, 0.8, 1.8, 3.7, 4.8, 
        "도메인 KPI 계산",
        [
            "■ 핵심 지표 SQL 마트 구축",
            "  - VPIP(자발적 참여율)",
            "  - PFR(선제공격률)",
            "  - AF(포스트플랍 공격성)",
            "  - Win Rate(최종 승률)",
            "",
            "■ 윈도우/집계 함수 최적화",
            "  - 932명 유저에 대한 지표를 중첩 스캔 없이 고속 산출하는 SQL 로직 구축"
        ]
    )
    
    add_card(
        slide4, 4.8, 1.8, 3.7, 4.8, 
        "통계적 가설 검정 (A/B Test)",
        [
            "■ 가설 정의",
            "  - 프리미엄 카드 보유 시 선공(A)과 단순 콜(B) 진입의 칩 획득량 차이 검증",
            "",
            "■ 분석 기법",
            "  - 독립표본 T-검정 수행",
            "",
            "■ 검정 결과 및 진단",
            "  - ★p-value = 0.6372로 가설 기각",
            "  - 귀무가설 수용 원인 규명: 수동적 콜 진입 집합(대조군)의 모수가 극소수여서 통계적 검정력이 부족했음을 정밀 진단함"
        ]
    )
    
    add_card(
        slide4, 8.8, 1.8, 3.7, 4.8, 
        "머신러닝 모델링 (ML)",
        [
            "■ K-Means 군집 분석",
            "  - VPIP, PFR 지표를 표준화하여 엘보우 기법 적용",
            "  - ★유저들을 4대 플레이 성향 코호트군으로 명확히 군집화",
            "",
            "■ Random Forest 예측 모델",
            "  - KPI 피처 기반 유저 흑자 여부 예측 Classifier 학습",
            "  - ★모델 성능: Accuracy 84%, ROC-AUC 0.91 달성"
        ]
    )
    
    # 슬라이드 4 노트 설정
    set_presenter_notes(slide4, """[슬라이드 4 발표 대본]
"세 번째로 통계적 가설 검정과 머신러닝 모델링 부분입니다.

먼저 포커 분석의 핵심인 VPIP(자발적 참여율), PFR(선제공격률) 등 도메인 KPI를 SQL 쿼리로 빠르게 가공해 데이터 마트를 구축했습니다.

그 후, 의사결정의 과학적 근거를 마련하고자 A/B 테스트 형식의 독립표본 T-검정을 설계했습니다. '좋은 카드를 잡았을 때 레이즈로 선제공격하며 참여한 집단(A)이 단순히 콜만 따고 소극적으로 진입한 집단(B)보다 수익이 높을 것이다'라는 가설을 세웠습니다.
검정 결과, p-value는 0.6372로 나와 가설이 기각되었으며 통계적으로 유의미한 차이를 보지 못했습니다.
저는 단순히 실패로 남겨두지 않고 왜 유의미하지 않은지를 추적 분석하였으며, 소극적으로 콜만 따고 들어간 B 집단의 표본 크기가 전체의 5% 미만으로 극도로 작아 통계적 검정력 자체가 부족했음을 데이터 구조상으로 밝혀냈습니다.

한편, K-Means 머신러닝 알고리즘을 적용해 932명의 유저들을 플레이 성향에 따라 TAG, LP 등 4대 플레이 성향 그룹으로 자동 군집화했습니다.
나아가 플레이어 지표를 넣어 이 플레이어가 최종적으로 돈을 딸지 잃을지를 예측하는 Random Forest 분류 모델을 연동하여 정확도 84%, ROC-AUC 0.91이라는 우수한 예측 성능을 확보했습니다."

[용어 리마인더]
- T-검정 (A/B Test): 두 그룹의 평균 차이가 진짜 의미 있는 차이인지, 우연히 우점한 것인지 판정하는 도구.
- p-value (유의확률): 우연히 그렇게 다를 확률. 0.05 미만이어야 우연이 아닌 유의미한 차이로 인정.
- K-Means 군집 분석: 인공지능이 유저를 자동으로 비슷한 성격끼리 묶어주는 알고리즘.
- Random Forest: 스무고개 놀이를 하듯 다수결의 질문 트리를 모아 최종 판단(돈 딸지 잃을지)을 내리는 모델.""")

    # =========================================================================
    # SLIDE 5: Actionable Outputs (분석 결과의 서비스화)
    # =========================================================================
    slide5 = create_slide_with_bg(prs)
    add_header(slide5, "04. 분석 결과의 서비스화: Streamlit 시뮬레이터", "PLATFORM SERVICE & SIMULATOR")
    
    add_card(
        slide5, 0.8, 1.8, 5.6, 4.8, 
        "플랫폼 아키텍처 및 대시보드 구축",
        [
            "■ Tableau 스타일 다차원 글로벌 필터",
            "  - 사이드바에서 토너먼트 종류, BB 슬라이더, 플레이어 포지션을 조절하면 실시간으로 연동된 SQL 쿼리가 파라미터 바인딩을 통해 전체 분석 차트를 즉시 동적 렌더링",
            "",
            "■ 오픈 웹 상용화 배포",
            "  - Streamlit Cloud 서비스를 이용해 외부에서도 포트폴리오를 조회하고 검증할 수 있도록 완전 공개 배포 완료",
            "  - 반응형 CSS 및 다국어(한국어/영어) 전환 라우팅을 구축하여 글로벌 유저 대응 역량 확보"
        ]
    )
    
    add_card(
        slide5, 6.9, 1.8, 5.6, 4.8, 
        "플레이 스타일 시뮬레이션 기능",
        [
            "■ 1. 실시간 ETL 적재 시뮬레이션 데모",
            "  - 유저가 로그 업로드 시 파일 파싱과 DB 적재 상황을 진행바로 시각화하며 Hero의 수익 곡선 실시간 렌더링",
            "",
            "■ 2. 플레이 스타일 샌드박스 플레이그라운드",
            "  - 유저가 VPIP, PFR, AF 슬라이더를 조정하여 플레이 패턴 설정 시 군집 성향 판정",
            "  - ★DB 내 가장 유사한 실제 유저 3명을 매칭하고 이들의 30세션 누적 자산 추이 시계열 그래프를 대조하여, '지표 개선 시 자산 파산(이탈) 확률 방어 여부'를 모의 실험 가능하게 구현"
        ]
    )
    
    # 슬라이드 5 노트 설정
    set_presenter_notes(slide5, """[슬라이드 5 발표 대본]
"네 번째로 분석 결과를 실제 서비스로 전환한 내용입니다.

저는 분석 결과를 단순한 정적 보고서로 방치하지 않고, 사용자가 직접 수치를 조작하며 상호작용할 수 있는 인터랙티브 Streamlit 웹 애플리케이션으로 배포하였습니다.

화면 왼쪽에는 토너먼트 유형이나 블라인드 단계를 조절할 수 있는 다차원 글로벌 필터 시스템을 구축하여, 유저가 마우스로 필터를 조절하는 즉시 백엔드에서 데이터베이스 쿼리를 동적으로 날려 플랫폼 차트 전체가 유기적으로 연동되도록 하였습니다.

특히, 핵심 성과물인 '플레이 스타일 샌드박스 시뮬레이터'는 유저가 본인의 VPIP, PFR, AF 수치를 입력하면, 머신러닝의 유클리드 거리 연산을 통해 DB에 저장된 932명 유저 중 성격이 가장 가까운 3명의 실제 플레이어를 실시간 추적 매칭합니다.
그리고 그들의 30세션 누적 자산(칩스) 시계열 추이 곡선을 화면에 대조하여 뿌려줍니다.
이를 통해 유저는 '내가 베팅 습관을 고치지 않으면 몇 게임 뒤에 파산하여 이탈(Churn)할 확률이 높겠구나' 혹은 '안정적으로 성장하겠구나'를 스스로 시뮬레이션할 수 있습니다."

[용어 리마인더]
- 유클리드 거리: 지표 공간 내 두 점 간의 직선거리. 여기선 '나와 가장 성격이 비슷한 3명의 유저'를 찾을 때 쓰는 거리 계산 공식.
- Streamlit(스트림릿): 파이썬 코드를 즉각 멋진 대화형 웹사이트로 변형해 주는 라이브러리.""")

    # =========================================================================
    # SLIDE 6: Successes & Improvements (성과 및 개선점)
    # =========================================================================
    slide6 = create_slide_with_bg(prs)
    add_header(slide6, "05. 성과 및 향후 개선점", "SUCCESS & ROADMAP")
    
    add_card(
        slide6, 0.8, 1.8, 5.6, 4.8, 
        "분석 과정에서 잘된 점 (Successes)",
        [
            "■ 1. 비정형 로그의 관계형 정형화 성공",
            "  - 의미 파악이 어렵던 텍스트 게임 이벤트를 Fact/Dimension 개념의 스타 스키마로 설계하여 다차원 글로벌 분석 쿼리 응답 성능 극대화",
            "",
            "■ 2. 통계와 머신러닝의 유기적 융합",
            "  - 단순 시각화를 넘어 귀무가설 검정을 통한 원인 진단, K-Means 성향 군집 분류, RF 예측으로 분석의 깊이 확보",
            "",
            "■ 3. 분석 결과의 MaaS화 (Service화)",
            "  - 정적 PDF/PPT 분석 보고서에 그치지 않고, 사용자가 직접 수치를 변경해보고 매칭 예측을 해보는 시뮬레이터 플랫폼을 제공하여 실질적 행동 개선 의사결정에 기여"
        ]
    )
    
    add_card(
        slide6, 6.9, 1.8, 5.6, 4.8, 
        "개선 필요 사항 및 향후 고도화 로드맵",
        [
            "■ 1. 통계적 표본 불균형 보완 전략 수립",
            "  - A/B 테스트에서 기각 원인이 된 대조군 표본 수 부족 해결",
            "  - *대안*: ★성향점수 매칭(PSM) 기법을 통한 통제군 형성 및 ★부트스트랩(Bootstrap) 리샘플링을 통한 신뢰구간 복원 적용 계획",
            "",
            "■ 2. 대규모 로그 스트리밍을 위한 확장형 파이프라인",
            "  - 단일 RDBMS의 동시성 락 및 용량 병목 극복",
            "  - *대안*: 실시간 수집을 위한 ★Apache Kafka 큐 및 실시간 처리를 위한 ★Apache Spark 분산 파이프라인 설계, 클라우드 DW(BigQuery)로의 적재 이전 로드맵 구상"
        ]
    )
    
    # 슬라이드 6 노트 설정
    set_presenter_notes(slide6, """[슬라이드 6 발표 대본]
"마지막으로 본 프로젝트의 성과와 향후 개선점에 대해 말씀드리겠습니다.

분석 과정에서의 대표적인 성과는 첫째, 텍스트 형태의 비정형 로그를 표준 RDBMS 스타 스키마로 적재하여 다차원 글로벌 조회 쿼리 성능을 대폭 향상하고 데이터웨어하우스 기반을 확립한 점입니다. 둘째, 단방향 통계 분석에 그치지 않고 유저의 실질적인 습관 교정과 의사결정을 돕는 인터랙티브 웹 시뮬레이터 플랫폼을 개발하고 완전히 오픈 상용 배포했다는 점입니다.

반면, 분석 과정에서의 통계적 한계와 아키텍처 확장성에 대한 개선 방향도 명확히 정의했습니다.

첫째, A/B 테스트에서 대조군의 모수가 적어 기각된 문제는 향후 성향점수 매칭(PSM) 기법과 부트스트랩 리샘플링 기술을 활용해 표본 편향을 수학적으로 완전히 극복할 계획입니다.

둘째, 향후 수천만 건 이상의 트래픽이 몰리는 대형 플랫폼 환경을 감안하면 단일 DBMS 아키텍처는 동시성 락으로 인해 병목이 생길 수 있습니다. 이를 고도화하기 위해 실시간 데이터 수집 큐인 Apache Kafka와 대용량 분산 처리를 위한 Apache Spark를 도입하고, 최종 저장소를 구글 클라우드의 BigQuery 데이터웨어하우스로 마이그레이션하는 빅데이터 파이프라인 로드맵을 수립하고 연구하고 있습니다.

데이터의 정제 수집부터 아키텍처 구성, 분석 임팩트 증명 및 최종 웹앱 서비스 배포까지 데이터의 흐름 전체를 주도해 본 경험을 바탕으로, SK Planet에서 실질적인 비즈니스 가치를 만들어 내는 기여자가 되겠습니다. 이상으로 발표를 마치겠습니다. 감사합니다."

[용어 리마인더]
- 성향점수 매칭 (PSM): A/B 테스트에서 두 집단의 조건 불균형을 잡기 위해 수학적으로 조건이 비슷한 짝을 찾아 맞춰주는 기법.
- 부트스트랩 (Bootstrap): 표본이 적을 때 데이터를 수만 번 복원 추출하여 가상 분포를 얻는 신뢰성 확보 통계법.
- Kafka & Spark: 실시간으로 쏟아지는 수백만 건의 대규모 트래픽 데이터를 처리해 주는 글로벌 표준 데이터 스택.""")

    # PPT 저장
    ppt_path = "poker_analytics_presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved successfully to: {os.path.abspath(ppt_path)}")

if __name__ == "__main__":
    main()
